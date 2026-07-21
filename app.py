# ============================================================================
# Sistema Multi-Agente SEI — App Streamlit
# Challenge AluraAgente — Oracle ONE AI for Tech
# Consolidado 18/07/2026 desde: dia3_rag_subagentes_v2.py, dia4_agente_pm_v2.py,
# sub_agentes_sql_final_v4.py (versiones vigentes, ver SKILL_04_CHALLENGE v4.1)
# ============================================================================
import os
import re
import glob
import time
from dataclasses import dataclass
from typing import List, Optional, TypedDict, Literal

import numpy as np
import streamlit as st
import faiss
from sentence_transformers import SentenceTransformer
from docx import Document as DocxDocument
import openpyxl
from pypdf import PdfReader
from pptx import Presentation

from sqlalchemy import create_engine, text
from langchain_community.utilities import SQLDatabase
from langchain_core.prompts import PromptTemplate, ChatPromptTemplate
from langchain_core.output_parsers import StrOutputParser
from langchain_core.messages import SystemMessage, HumanMessage
from langchain_groq import ChatGroq
from langgraph.graph import StateGraph, END
from pydantic import BaseModel, Field
from groq import RateLimitError

st.set_page_config(page_title="Sistema Multi-Agente SEI", page_icon="🏗️", layout="centered")

# ============================================================================
# CONFIGURACIÓN — variables desde Streamlit Secrets (Settings → Secrets)
# ============================================================================
GROQ_API_KEY = st.secrets["GROQ_API_KEY"]
DATABASE_URL = st.secrets["DATABASE_URL"]
os.environ["GROQ_API_KEY"] = GROQ_API_KEY

DATA_DIR = os.path.join(os.path.dirname(__file__), "data")
FECHA_REFERENCIA_DEMO = "2026-07-13"  # fecha fija del dataset de demo

# ============================================================================
# RECURSOS CACHEADOS — se cargan una sola vez por instancia de la app
# ============================================================================

@st.cache_resource(show_spinner="Conectando a la base de datos...")
def get_engine():
    return create_engine(DATABASE_URL)

@st.cache_resource(show_spinner="Cargando modelo de lenguaje...")
def get_llm():
    return ChatGroq(model="openai/gpt-oss-120b", temperature=0)


# ---- RAG: extracción, metadata, chunking (idéntico a dia3_rag_subagentes_v2.py) ----

def leer_md_txt(path: str) -> str:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()

def leer_docx(path: str) -> str:
    doc = DocxDocument(path)
    partes = [p.text for p in doc.paragraphs if p.text.strip()]
    for tabla in doc.tables:
        for fila in tabla.rows:
            celdas = [c.text.strip() for c in fila.cells if c.text.strip()]
            if celdas:
                partes.append(" | ".join(celdas))
    return "\n".join(partes)

def leer_xlsx(path: str) -> str:
    wb = openpyxl.load_workbook(path, data_only=True)
    partes = []
    for hoja in wb.sheetnames:
        ws = wb[hoja]
        for fila in ws.iter_rows(values_only=True):
            valores = [str(v).strip() for v in fila if v is not None and str(v).strip()]
            if valores:
                partes.append(" | ".join(valores))
    return "\n".join(partes)

def leer_pdf(path: str) -> str:
    reader = PdfReader(path)
    partes = []
    for pagina in reader.pages:
        texto = pagina.extract_text()
        if texto:
            partes.append(texto)
    return "\n".join(partes)

def leer_pptx(path: str) -> str:
    prs = Presentation(path)
    partes = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                partes.append(shape.text_frame.text)
            if shape.has_table:
                for fila in shape.table.rows:
                    celdas = [c.text.strip() for c in fila.cells if c.text.strip()]
                    if celdas:
                        partes.append(" | ".join(celdas))
    return "\n".join(partes)

def extraer_texto(path: str) -> str:
    ext = path.lower().rsplit(".", 1)[-1]
    try:
        if ext in ("md", "txt"):
            return leer_md_txt(path)
        if ext == "docx":
            return leer_docx(path)
        if ext == "xlsx":
            return leer_xlsx(path)
        if ext == "pdf":
            return leer_pdf(path)
        if ext == "pptx":
            return leer_pptx(path)
    except Exception:
        pass
    return ""

PATRON_CODIGO = re.compile(r"((?:SE0[12]|SEI)-[\w]+-[A-Z]-[A-Z]+-\d+)_([A-Z0-9]+)")

def parsear_metadata(path: str, proyecto: str) -> dict:
    nombre = os.path.basename(path)
    m = PATRON_CODIGO.search(nombre)
    codigo = m.group(1) if m else nombre
    revision = m.group(2) if m else ""
    vigente = "/SUP/" not in path.replace("\\", "/")

    ruta_norm = path.replace("\\", "/")
    if proyecto == "GENERAL":
        tipo_documento = "INSTITUCIONAL"
    elif "MINUTAS REUNION" in ruta_norm:
        tipo_documento = "MINUTA"
    elif "00-RFI" in ruta_norm and "SDI" in nombre.upper():
        tipo_documento = "SDI"
    elif "00-RFI" in ruta_norm and "RFI" in nombre.upper():
        tipo_documento = "STATUS_RFI"
    elif "00-TTR" in ruta_norm:
        tipo_documento = "TRANSMITTAL"
    elif "00-DOC REC" in ruta_norm:
        tipo_documento = "CORREO"
    elif "03-INFO REFERENCIA" in ruta_norm:
        tipo_documento = "NORMA_REFERENCIA"
    else:
        tipo_documento = "TECNICO"

    disciplina = ""
    for carpeta in ["01-G", "02-M", "03-P", "04-A", "05-C", "06-E"]:
        if f"/{carpeta}/" in ruta_norm:
            disciplina = carpeta.split("-")[1]
            break
    if proyecto == "SE01" and tipo_documento == "TECNICO":
        disciplina = "E"

    return {"proyecto": proyecto, "ruta": path, "documento": codigo, "revision": revision,
            "vigente": vigente, "tipo_documento": tipo_documento, "disciplina": disciplina}

EXTENSIONES_SOPORTADAS = (".pdf", ".docx", ".xlsx", ".md", ".txt", ".pptx")

def es_conocimiento(path: str) -> bool:
    """Clasifica por CARPETA únicamente. La extensión se valida acá solo
    para descartar formatos no soportados (imágenes sueltas, etc.) — un PDF,
    un Word, un Excel o un PPT en la misma carpeta se tratan exactamente igual."""
    p = path.replace("\\", "/")
    if not p.lower().endswith(EXTENSIONES_SOPORTADAS):
        return False
    if "/GENERAL/" in p:
        return True
    if "03-INFO REFERENCIA" in p:
        return True
    if "02-ELABORADOS" in p and "00-MINUTAS REUNION" not in p and "00-RFI" not in p:
        return True
    return False

def es_comunicaciones(path: str) -> bool:
    """Idem es_conocimiento: clasifica por carpeta, no por extensión."""
    p = path.replace("\\", "/")
    if not p.lower().endswith(EXTENSIONES_SOPORTADAS):
        return False
    if "00-MINUTAS REUNION" in p:
        return True
    if "00-RFI" in p:
        return True
    if "00-DOC REC" in p:
        return True
    if "00-TTR" in p:
        return True
    return False

@dataclass
class Chunk:
    texto: str
    metadata: dict

def dividir_en_chunks(texto: str, tam: int = 700, solape: int = 100) -> List[str]:
    texto = texto.strip()
    if len(texto) <= tam:
        return [texto] if texto else []
    chunks, inicio = [], 0
    while inicio < len(texto):
        fin = min(inicio + tam, len(texto))
        chunks.append(texto[inicio:fin])
        inicio += tam - solape
    return chunks

def construir_corpus(rutas_proyectos: dict, filtro_fn) -> List[Chunk]:
    corpus = []
    for proyecto, raiz in rutas_proyectos.items():
        archivos = glob.glob(os.path.join(raiz, "**", "*"), recursive=True)
        archivos = [a for a in archivos if os.path.isfile(a) and filtro_fn(a)]
        for path in archivos:
            texto = extraer_texto(path)
            if not texto.strip():
                continue
            meta = parsear_metadata(path, proyecto)
            for i, frag in enumerate(dividir_en_chunks(texto)):
                m = dict(meta)
                m["chunk_id"] = i
                corpus.append(Chunk(texto=frag, metadata=m))
    return corpus

class IndiceFAISS:
    def __init__(self, chunks: List[Chunk], embedder):
        self.chunks = chunks
        self.embedder = embedder
        if not chunks:
            self.index = None
            return
        vecs = self._embeber([c.texto for c in chunks], "passage: ")
        self.index = faiss.IndexFlatIP(vecs.shape[1])
        self.index.add(vecs)

    def _embeber(self, textos: List[str], prefijo: str) -> np.ndarray:
        textos_prefijados = [f"{prefijo}{t}" for t in textos]
        vecs = self.embedder.encode(textos_prefijados, normalize_embeddings=True,
                                     show_progress_bar=False)
        return np.array(vecs, dtype="float32")

    def buscar(self, query: str, k: int = 4, filtro: Optional[dict] = None) -> List[Chunk]:
        if self.index is None:
            return []
        qvec = self._embeber([query], "query: ")
        k_busqueda = len(self.chunks)
        _, idxs = self.index.search(qvec, k_busqueda)
        resultados = []
        for i in idxs[0]:
            if i < 0:
                continue
            c = self.chunks[i]
            if filtro and not all(c.metadata.get(k2) == v2 for k2, v2 in filtro.items()):
                continue
            resultados.append(c)
            if len(resultados) >= k:
                break
        return resultados

def formatear_contexto(chunks: List[Chunk]) -> str:
    partes = []
    for c in chunks:
        m = c.metadata
        estado = "VIGENTE" if m["vigente"] else "SUPERADO"
        partes.append(f"[{m['documento']} rev.{m['revision']} | {m['tipo_documento']} | "
                       f"{m['proyecto']} | {estado}]\n{c.texto}")
    return "\n\n---\n\n".join(partes)

PROMPT_RAG = PromptTemplate.from_template(
    """Sos un asistente técnico de SEI Ingeniería S.R.L. Respondé la pregunta del
usuario ÚNICAMENTE con la información de los fragmentos de contexto de abajo.
Si el contexto no alcanza para responder, decilo explícitamente — no inventes datos.
Citá siempre el código de documento y revisión entre paréntesis cuando corresponda.

CONTEXTO:
{contexto}

PREGUNTA: {pregunta}

RESPUESTA:"""
)

@st.cache_resource(show_spinner="Indexando base de conocimiento (RAG)... puede tardar 1-2 min la primera vez")
def build_rag():
    embedder = SentenceTransformer("intfloat/multilingual-e5-small")
    rutas = {"SE01": os.path.join(DATA_DIR, "SE01"), "SE02": os.path.join(DATA_DIR, "SE02"),
             "GENERAL": os.path.join(DATA_DIR, "GENERAL")}
    corpus_conocimiento = construir_corpus(rutas, es_conocimiento)
    corpus_comunicaciones = construir_corpus(rutas, es_comunicaciones)
    idx_conocimiento = IndiceFAISS(corpus_conocimiento, embedder)
    idx_comunicaciones = IndiceFAISS(corpus_comunicaciones, embedder)
    return idx_conocimiento, idx_comunicaciones


# ============================================================================
# SUB-AGENTES SQL (idéntico a sub_agentes_sql_final_v4.py)
# ============================================================================
_PALABRAS_PROHIBIDAS = re.compile(
    r"\b(INSERT|UPDATE|DELETE|DROP|ALTER|TRUNCATE|CREATE|REPLACE|GRANT|REVOKE|EXEC|EXECUTE|CALL)\b",
    flags=re.IGNORECASE,
)

def _es_select_seguro(sql_query: str) -> bool:
    limpio = sql_query.strip()
    while True:
        antes = limpio
        limpio = re.sub(r"^\s*--[^\n]*\n?", "", limpio)
        limpio = re.sub(r"^\s*/\*.*?\*/", "", limpio, flags=re.S)
        limpio = limpio.strip()
        if limpio == antes:
            break
    if not re.match(r"^(SELECT|WITH)\b", limpio, flags=re.IGNORECASE):
        return False
    if _PALABRAS_PROHIBIDAS.search(limpio):
        return False
    if ";" in limpio:
        resto = limpio.split(";", 1)[1].strip()
        if resto:
            return False
    return True

_PROMPT_SQL_AVANCE = PromptTemplate.from_template("""
Eres un asistente que traduce preguntas en español sobre gestión de proyectos
de ingeniería a consultas SQL para MySQL.

Contexto importante:
- SPI (Schedule Performance Index) = columna `indice_desempeno_cronograma`
- CPI (Cost Performance Index) = columna `indice_desempeno_costos`
- Los proyectos se identifican por `codigo_proyecto` ('SE01', 'SE02') en la
  tabla `proyectos`, pero las tablas de gestión usan `id_proyecto` (FK numérica).
  Si la pregunta menciona 'SE01' o 'SE02', hacé JOIN con `proyectos` usando
  `codigo_proyecto` para resolver el `id_proyecto`.
- Solo hay 2 proyectos cargados: SE01 (id_proyecto=9) y SE02 (id_proyecto=10).
- Para cualquier expresión de tiempo relativo ("la semana pasada", "este mes",
  "hoy", "últimos N días"), NO uses CURDATE() ni NOW(). Usá como fecha de
  referencia fija el literal '{fecha_referencia}'.

Esquema disponible:
{table_info}

Pregunta: {pregunta}

Generá SOLO la consulta SQL para MySQL (sin explicación, sin markdown, sin punto y coma final).
""")

_PROMPT_RESPUESTA_AVANCE = PromptTemplate.from_template("""
Sos el sub-agente SQL-Avance del Sistema Multi-Agente SEI. Respondé la pregunta
del usuario en español, de forma directa y profesional (estilo ingeniero a
ingeniero), basándote ÚNICAMENTE en el resultado SQL que se te da. Si el
resultado está vacío, decilo explícitamente — nunca inventes datos.

Pregunta: {pregunta}
Resultado SQL: {resultado}

Respuesta:
""")

_PROMPT_TIMESHEET = PromptTemplate.from_template("""
Eres un asistente que traduce preguntas en español sobre horas trabajadas
(timesheet) de ingenieros a consultas SQL para MySQL.

Contexto importante:
- `registro_horas.descripcion_tarea` es el campo con el detalle real de la
  tarea (texto libre) — NO existe una columna llamada 'comentarios'.
- `registro_horas.horas_normales` es la cantidad de horas cargadas.
- Cada registro está vinculado a UN documento (`id_documento`) y UN tipo de
  tarea (`id_tipo_tarea`, tabla `tipos_tarea`).
- Para filtrar por proyecto (SE01/SE02), unir registro_horas -> documentos
  -> proyectos usando codigo_proyecto.
- Para filtrar por persona, unir con `empleados` (nombre, apellido).
- Si la pregunta pide "horas totales" en cierto tipo de documento (ej.
  "memorias de cálculo"), filtrar por `documentos.id_tipo_documento` = 2
  (MC) y usar SUM(horas_normales), agrupando por proyecto si corresponde.
- Para cualquier expresión de tiempo relativo, NO uses CURDATE() ni NOW().
  Usá como fecha de referencia fija el literal '{fecha_referencia}'.

Esquema disponible:
{table_info}

Pregunta: {pregunta}

Generá SOLO la consulta SQL para MySQL (sin explicación, sin markdown, sin punto y coma final).
""")

_PROMPT_RESPUESTA_TIMESHEET = PromptTemplate.from_template("""
Sos el sub-agente Timesheet del Sistema Multi-Agente SEI. Respondé la
pregunta del usuario en español, de forma directa (estilo ingeniero a
ingeniero), basándote ÚNICAMENTE en el resultado SQL que se te da. Si pide
un desglose (por tipo de tarea, por documento), presentalo claro. Si el
resultado está vacío, decilo explícitamente — nunca inventes datos.

Pregunta: {pregunta}
Resultado SQL: {resultado}

Respuesta:
""")

def sql_avance(pregunta: str) -> dict:
    engine, llm = get_engine(), get_llm()
    db = SQLDatabase(engine, include_tables=["gestion_proyectos", "riesgos_proyecto",
                                              "hitos_proyecto", "cambios_alcance", "proyectos"])
    chain = _PROMPT_SQL_AVANCE | llm | StrOutputParser()
    resp_chain = _PROMPT_RESPUESTA_AVANCE | llm | StrOutputParser()
    sql_query = ""
    try:
        sql_query = chain.invoke({"table_info": db.get_table_info(), "pregunta": pregunta,
                                   "fecha_referencia": FECHA_REFERENCIA_DEMO})
        sql_query = sql_query.replace("```sql", "").replace("```", "").strip()
        if not _es_select_seguro(sql_query):
            return {"respuesta": "No pude procesar esa consulta: la traducción a SQL generada "
                                  "no es una lectura segura (SELECT).", "sql_generado": sql_query,
                    "resultado_crudo": [], "sub_agente": "sql_avance",
                    "error": "guardrail: consulta no es SELECT"}
        with engine.connect() as conn:
            resultado = conn.execute(text(sql_query)).fetchall()
        respuesta = resp_chain.invoke({"pregunta": pregunta, "resultado": str(resultado)})
        return {"respuesta": respuesta, "sql_generado": sql_query,
                "resultado_crudo": [tuple(r) for r in resultado], "sub_agente": "sql_avance", "error": None}
    except Exception as e:
        return {"respuesta": f"No pude responder esa consulta de avance del proyecto. Error: {e}",
                "sql_generado": sql_query, "resultado_crudo": [], "sub_agente": "sql_avance", "error": str(e)}

def timesheet(pregunta: str) -> dict:
    engine, llm = get_engine(), get_llm()
    db = SQLDatabase(engine, include_tables=["registro_horas", "empleados", "documentos",
                                              "tipos_tarea", "proyectos"])
    chain = _PROMPT_TIMESHEET | llm | StrOutputParser()
    resp_chain = _PROMPT_RESPUESTA_TIMESHEET | llm | StrOutputParser()
    sql_query = ""
    try:
        sql_query = chain.invoke({"table_info": db.get_table_info(), "pregunta": pregunta,
                                   "fecha_referencia": FECHA_REFERENCIA_DEMO})
        sql_query = sql_query.replace("```sql", "").replace("```", "").strip()
        if not _es_select_seguro(sql_query):
            return {"respuesta": "No pude procesar esa consulta: la traducción a SQL generada "
                                  "no es una lectura segura (SELECT).", "sql_generado": sql_query,
                    "resultado_crudo": [], "sub_agente": "timesheet",
                    "error": "guardrail: consulta no es SELECT"}
        with engine.connect() as conn:
            resultado = conn.execute(text(sql_query)).fetchall()
        respuesta = resp_chain.invoke({"pregunta": pregunta, "resultado": str(resultado)})
        return {"respuesta": respuesta, "sql_generado": sql_query,
                "resultado_crudo": [tuple(r) for r in resultado], "sub_agente": "timesheet", "error": None}
    except Exception as e:
        return {"respuesta": f"No pude responder esa consulta de horas. Error: {e}",
                "sql_generado": sql_query, "resultado_crudo": [], "sub_agente": "timesheet", "error": str(e)}

def rag_conocimiento(pregunta: str, proyecto: Optional[str] = None,
                      solo_vigentes: bool = True, k: int = 4) -> dict:
    idx_conocimiento, _ = build_rag()
    llm = get_llm()
    filtro = {}
    if proyecto:
        filtro["proyecto"] = proyecto
    if solo_vigentes:
        filtro["vigente"] = True
    chunks = idx_conocimiento.buscar(pregunta, k=k, filtro=filtro or None)
    if not chunks:
        return {"respuesta": "No encontré documentación técnica relevante para esa consulta.",
                "citaciones": [], "documentos_encontrados": False}
    contexto = formatear_contexto(chunks)
    cadena = PROMPT_RAG | llm | StrOutputParser()
    respuesta = cadena.invoke({"contexto": contexto, "pregunta": pregunta})
    return {"respuesta": respuesta, "citaciones": chunks, "documentos_encontrados": True}

def rag_comunicaciones(pregunta: str, proyecto: Optional[str] = None,
                        solo_vigentes: bool = True, k: int = 4) -> dict:
    _, idx_comunicaciones = build_rag()
    llm = get_llm()
    filtro = {}
    if proyecto:
        filtro["proyecto"] = proyecto
    if solo_vigentes:
        filtro["vigente"] = True
    chunks = idx_comunicaciones.buscar(pregunta, k=k, filtro=filtro or None)
    if not chunks:
        return {"respuesta": "No encontré comunicaciones (minutas/SDI/transmittals) relevantes.",
                "citaciones": [], "documentos_encontrados": False}
    contexto = formatear_contexto(chunks)
    cadena = PROMPT_RAG | llm | StrOutputParser()
    respuesta = cadena.invoke({"contexto": contexto, "pregunta": pregunta})
    return {"respuesta": respuesta, "citaciones": chunks, "documentos_encontrados": True}


# ============================================================================
# AGENTE PM — ORQUESTADOR (idéntico a dia4_agente_pm_v2.py)
# ============================================================================
def _adaptar_respuesta_sql(resultado) -> dict:
    if isinstance(resultado, dict) and "respuesta" in resultado:
        return resultado
    return {"respuesta": str(resultado), "citaciones": [], "documentos_encontrados": True}

class TriajeOut(BaseModel):
    agentes: List[Literal["SQL_AVANCE", "TIMESHEET", "RAG_CONOCIMIENTO", "RAG_COMUNICACIONES"]] = Field(
        description="Uno o más sub-agentes necesarios para responder la pregunta.")
    proyecto: Optional[Literal["SE01", "SE02"]] = Field(default=None)
    razonamiento: str = Field(description="Breve justificación de la clasificación (1 línea).")

PROMPT_TRIAJE = """Sos el clasificador de un sistema multi-agente de ingeniería de
SEI Ingeniería S.R.L. Dada una pregunta en lenguaje natural, decidí qué sub-agente(s)
deben responderla. Reglas de asignación:

- SQL_AVANCE: preguntas sobre avance de proyecto, SPI/CPI, riesgos, hitos,
  cambios de alcance — datos de gestión estructurados en base de datos.
- TIMESHEET: preguntas sobre horas cargadas, quién trabajó en qué tarea/documento,
  totales de horas por categoría o persona.
- RAG_CONOCIMIENTO: preguntas sobre contenido técnico de documentos de ingeniería
  vigentes (memorias de cálculo, listados de cargas, planos, normas aplicadas,
  lecciones aprendidas de diseño).
- RAG_COMUNICACIONES: preguntas sobre qué se acordó, consultó o respondió en
  minutas de reunión, SDI/RFI, o transmittals con el cliente.

Si la pregunta necesita combinar datos de gestión/lecciones aprendidas (BD) CON
el contenido de un documento técnico, incluí SQL_AVANCE (o TIMESHEET) Y
RAG_CONOCIMIENTO juntos en la lista — no elijas solo uno si la pregunta pide ambos.

Si la pregunta menciona explícitamente "SE01" o "SE02", o un cliente
(Minera NOA Litio SA = SE01, Litio Andino SA = SE02), completá el campo proyecto.
Si no se menciona ningún proyecto, dejalo en null (no asumas)."""

class AgentState(TypedDict):
    pregunta: str
    triaje: dict
    resultados: List[dict]
    respuesta_final: str

PROMPT_SINTESIS = ChatPromptTemplate.from_messages([
    ("system", """Sos el Agente PM de SEI Ingeniería S.R.L. Recibiste los resultados de uno o
más sub-agentes especializados para responder la pregunta de un colaborador.
Combiná esos resultados en UNA respuesta clara y directa en español, ingeniero
a ingeniero, sin verborragia. Si un sub-agente no encontró información, decilo
en vez de omitirlo en silencio. No inventes nada que no esté en los resultados."""),
    ("human", """PREGUNTA ORIGINAL: {pregunta}

RESULTADOS DE LOS SUB-AGENTES:
{resultados}

RESPUESTA FINAL (síntesis):"""),
])

@st.cache_resource(show_spinner=False)
def build_agente_pm():
    llm = get_llm()
    llm_triaje = llm.with_structured_output(TriajeOut)
    cadena_sintesis = PROMPT_SINTESIS | llm

    def triaje(pregunta: str) -> dict:
        salida: TriajeOut = llm_triaje.invoke([
            SystemMessage(content=PROMPT_TRIAJE), HumanMessage(content=pregunta)])
        return salida.model_dump()

    def nodo_triaje(state: AgentState) -> AgentState:
        return {"triaje": triaje(state["pregunta"])}

    def nodo_ejecutar_subagentes(state: AgentState) -> AgentState:
        pregunta = state["pregunta"]
        proyecto = state["triaje"].get("proyecto")
        resultados = []
        for agente in state["triaje"]["agentes"]:
            if agente == "SQL_AVANCE":
                r = _adaptar_respuesta_sql(sql_avance(pregunta))
            elif agente == "TIMESHEET":
                r = _adaptar_respuesta_sql(timesheet(pregunta))
            elif agente == "RAG_CONOCIMIENTO":
                r = rag_conocimiento(pregunta, proyecto=proyecto)
            elif agente == "RAG_COMUNICACIONES":
                r = rag_comunicaciones(pregunta, proyecto=proyecto)
            else:
                continue
            r["agente"] = agente
            resultados.append(r)
        return {"resultados": resultados}

    def nodo_sintesis(state: AgentState) -> AgentState:
        resultados_fmt = "\n\n".join(f"[{r['agente']}]\n{r['respuesta']}" for r in state["resultados"])
        respuesta = cadena_sintesis.invoke({"pregunta": state["pregunta"], "resultados": resultados_fmt}).content
        return {"respuesta_final": respuesta}

    workflow = StateGraph(AgentState)
    workflow.add_node("triaje", nodo_triaje)
    workflow.add_node("ejecutar_subagentes", nodo_ejecutar_subagentes)
    workflow.add_node("sintesis", nodo_sintesis)
    workflow.set_entry_point("triaje")
    workflow.add_edge("triaje", "ejecutar_subagentes")
    workflow.add_edge("ejecutar_subagentes", "sintesis")
    workflow.add_edge("sintesis", END)
    return workflow.compile()

def preguntar_pm(pregunta: str, max_intentos: int = 4) -> dict:
    agente_pm = build_agente_pm()
    for intento in range(max_intentos):
        try:
            resultado = agente_pm.invoke({"pregunta": pregunta})
            return {"pregunta": pregunta, "agentes_consultados": resultado["triaje"]["agentes"],
                    "proyecto": resultado["triaje"].get("proyecto"),
                    "respuesta": resultado["respuesta_final"], "detalle_por_agente": resultado["resultados"]}
        except RateLimitError:
            time.sleep(15 * (intento + 1))
    raise RuntimeError("Límite de uso de Groq alcanzado. Esperá un minuto e intentá de nuevo.")


# ============================================================================
# INTERFAZ STREAMLIT
# ============================================================================
st.title("🏗️ Sistema Multi-Agente SEI")
st.caption("Challenge AluraAgente — Consulta en lenguaje natural sobre avance, "
           "horas, documentación técnica y comunicaciones de proyectos SE01/SE02.")

with st.expander("ℹ️ Ejemplos de preguntas"):
    st.markdown("""
- ¿Cuál es el SPI actual del proyecto SE01?
- ¿Qué riesgos abiertos tiene el proyecto SE02?
- ¿Qué dice el listado de cargas eléctricas de SE01 sobre el sector Evaporación?
- ¿Qué se acordó en la minuta de SE01 sobre el esquema de puesta a tierra?
- ¿Qué hizo Lucía Fernández en el proyecto SE02 la semana pasada?
- ¿Cuántas horas totales se cargaron a memorias de cálculo entre SE01 y SE02?
""")

pregunta = st.text_input("Escribí tu pregunta:", placeholder="Ej: ¿Cuál es el SPI actual del proyecto SE01?")

if st.button("Consultar", type="primary") and pregunta.strip():
    with st.spinner("Consultando al Agente PM..."):
        try:
            r = preguntar_pm(pregunta)
            st.markdown("### Respuesta")
            st.markdown(r["respuesta"])
            with st.expander("Detalle técnico (agentes consultados)"):
                st.write(f"**Agentes:** {', '.join(r['agentes_consultados'])}")
                if r["proyecto"]:
                    st.write(f"**Proyecto detectado:** {r['proyecto']}")
                for d in r["detalle_por_agente"]:
                    st.markdown(f"**[{d.get('agente')}]**")
                    if "sql_generado" in d and d["sql_generado"]:
                        st.code(d["sql_generado"], language="sql")
        except Exception as e:
            st.error(f"Ocurrió un error al procesar la consulta: {e}")

st.divider()
st.caption("SEI Ingeniería S.R.L. — Proyectos SE01 (Ampliación Planta Evaporación Salar Norte) "
           "y SE02 (Ampliación Almacenes Repuestos). Criterio anti-alucinación aplicado en los 4 "
           "sub-agentes: si un dato no está en la fuente consultada, se informa explícitamente.")
